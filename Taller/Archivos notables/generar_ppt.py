import collections.abc
import pptx
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_modern_ppt():
    prs = Presentation()
    
    # Configurar tamaño panorámico (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Paleta de colores minimalista y profesional
    COLOR_BG = RGBColor(248, 249, 250)      # Gris muy claro/blanco
    COLOR_PRIMARY = RGBColor(43, 87, 154)   # Azul corporativo elegante
    COLOR_ACCENT = RGBColor(230, 81, 0)     # Naranja oscuro para destacar
    COLOR_TEXT = RGBColor(33, 37, 41)       # Gris oscuro casi negro
    COLOR_MUTED = RGBColor(108, 117, 125)   # Gris medio para texto secundario
    COLOR_LIGHT = RGBColor(255, 255, 255)   # Blanco
    
    def apply_slide_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title(slide, text, top=0.5, font_size=40, color=COLOR_PRIMARY):
        title = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.7), Inches(1))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        return title

    def add_subtitle(slide, text, top=1.3, font_size=20):
        subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.7), Inches(0.5))
        tf = subtitle.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_MUTED
        p.font.name = 'Calibri'
        return subtitle
        
    def add_text(slide, text, left, top, width, height, font_size=18, bold=False, color=COLOR_TEXT):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        return box
        
    def add_kpi(slide, label, value, left, top, is_highlight=False):
        # Valor grande
        val_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(3), Inches(1))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT if is_highlight else COLOR_PRIMARY
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
        
        # Etiqueta abajo
        lbl_box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.9), Inches(3), Inches(0.5))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_MUTED
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER

    def add_insight_box(slide, text, top=6.2):
        # Rectángulo de fondo
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(0.8), Inches(top), Inches(11.7), Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_PRIMARY
        shape.line.color.rgb = COLOR_PRIMARY
        
        # Texto
        tf = shape.text_frame
        tf.vertical_anchor = pptx.enum.text.MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"Insight clave: {text}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_LIGHT
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER

    blank_layout = prs.slide_layouts[6]

    # --- SLIDE 1: PORTADA ---
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide, COLOR_BG)
    
    # Decoración geométrica simple
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.color.rgb = COLOR_PRIMARY

    add_title(slide, "Supermarket Sales Analysis", top=2.5, font_size=54)
    add_subtitle(slide, "Hallazgos de datos y estrategia para campañas de marketing", top=3.5, font_size=24)
    
    # KPIs rápidos en portada
    add_kpi(slide, "Transacciones", "1,000", 1.5, 5)
    add_kpi(slide, "Ventas Totales", "$322.9K", 5.16, 5)
    add_kpi(slide, "Ticket Promedio", "$323", 8.8, 5)

    # --- SLIDE 2: CONTEXTO ---
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide, COLOR_BG)
    add_title(slide, "Contexto del Negocio")
    add_subtitle(slide, "3 sucursales, 3 perfiles de ciudades distintas en Myanmar")
    
    # 3 columnas para las sucursales
    add_text(slide, "Yangon (Sucursal A)", 1, 2.5, 3.5, 0.5, 22, True)
    add_text(slide, "• Capital comercial\n• 340 transacciones\n• Competitivo, ticket bajo", 1, 3, 3.5, 2)
    
    add_text(slide, "Mandalay (Sucursal B)", 5, 2.5, 3.5, 0.5, 22, True)
    add_text(slide, "• Centro cultural\n• 332 transacciones\n• Foco en bienestar y salud", 5, 3, 3.5, 2)
    
    add_text(slide, "Naypyitaw (Sucursal C)", 9, 2.5, 3.5, 0.5, 22, True)
    add_text(slide, "• Capital política\n• 328 transacciones\n• Menos volumen, ticket más alto", 9, 3, 3.5, 2)

    add_insight_box(slide, "El objetivo es maximizar ingresos diseñando 3 campañas mensuales adaptadas al consumo local.")

    # --- SLIDE 3: GÉNERO ---
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide, COLOR_BG)
    add_title(slide, "Comportamiento por Género")
    add_subtitle(slide, "Las mujeres generan mayor valor, pero los hombres dominan nichos específicos")
    
    add_kpi(slide, "Ticket Promedio Mujeres", "$335", 2, 2.5, True)
    add_kpi(slide, "Ticket Promedio Hombres", "$311", 7.5, 2.5)
    
    add_text(slide, "• Las mujeres representan el 52% de las ventas totales ($167.8K vs $155K).", 1, 4.5, 11, 0.5, 20)
    add_text(slide, "• Sin embargo, en la categoría 'Health & Beauty', los hombres gastan un 65% más que las mujeres.", 1, 5, 11, 0.5, 20)
    
    add_insight_box(slide, "Oportunidad para romper estereotipos: promocionar productos de cuidado personal dirigidos a hombres.")

    # --- SLIDE 4: HORARIOS ---
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide, COLOR_BG)
    add_title(slide, "El ritmo de las ventas diarias")
    add_subtitle(slide, "Dos momentos clave concentran el tráfico")
    
    add_text(slide, "Pico del Almuerzo (13:00 hs)", 2, 2.5, 4, 0.5, 24, True, COLOR_PRIMARY)
    add_text(slide, "Primer momento de alto tráfico del día. Ideal para compras rápidas de alimentos y bebidas.", 2, 3.1, 4, 1.5, 18)
    
    add_text(slide, "Pico Post-Trabajo (19:00 hs)", 7.3, 2.5, 4, 0.5, 24, True, COLOR_ACCENT)
    add_text(slide, "El momento de mayor volumen de ventas absoluto. Compras más grandes y planificadas.", 7.3, 3.1, 4, 1.5, 18)
    
    add_text(slide, "Horas valle: 16:00, 17:00 y 20:00. Necesitan incentivos de tráfico.", 1, 5, 11, 0.5, 20, True, COLOR_MUTED)
    
    add_insight_box(slide, "Programar activaciones en sucursal y ofertas relámpago durante los picos de las 13h y 19h.")

    # --- SLIDE 5: PRODUCTOS ---
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide, COLOR_BG)
    add_title(slide, "Desempeño del Portafolio")
    add_subtitle(slide, "Alimentos lidera los ingresos; Fidelización aumenta el ticket")
    
    # Top y Bottom
    add_text(slide, "Categoría Estrella", 2, 2.5, 4, 0.5, 22, True)
    add_text(slide, "Food & Beverages\n$56.1K Ventas\nLiderada por mujeres y por la sucursal de Naypyitaw.", 2, 3, 4, 1.5)
    
    add_text(slide, "Categoría a Reforzar", 7.3, 2.5, 4, 0.5, 22, True)
    add_text(slide, "Health & Beauty\n$49.1K Ventas\nMenos transacciones generales, pero ticket alto en hombres.", 7.3, 3, 4, 1.5)
    
    add_text(slide, "El rol de la Membresía: Los clientes 'Member' tienen un ticket promedio $9.67 más alto que los 'Normal'.", 1, 5, 11, 0.5, 20, True, COLOR_PRIMARY)
    
    add_insight_box(slide, "Usar las categorías líderes como gancho para vender las más débiles mediante promociones cruzadas.")

    # --- SLIDE 6: ESTRATEGIA ---
    slide = prs.slides.add_slide(blank_layout)
    apply_slide_bg(slide, COLOR_BG)
    add_title(slide, "Propuesta: 3 Campañas, 3 Meses")
    add_subtitle(slide, "Acciones accionables basadas en la data")
    
    add_text(slide, "Mes 1: 'Sabores y Familia'", 1, 2.5, 3.5, 0.5, 22, True, COLOR_ACCENT)
    add_text(slide, "Foco: Food & Beverages\nTarget: Mujeres en Naypyitaw\nTáctica: Degustaciones a las 13h y 19h para impulsar el mes de más ventas.", 1, 3, 3.5, 2, 18)
    
    add_text(slide, "Mes 2: 'Bienestar para Él'", 5, 2.5, 3.5, 0.5, 22, True, COLOR_PRIMARY)
    add_text(slide, "Foco: Health & Beauty\nTarget: Hombres en Mandalay\nTáctica: Aprovechar el nicho masculino para levantar febrero (mes más bajo).", 5, 3, 3.5, 2, 18)
    
    add_text(slide, "Mes 3: 'Tu Casa a tu Estilo'", 9, 2.5, 3.5, 0.5, 22, True, COLOR_ACCENT)
    add_text(slide, "Foco: Home & Lifestyle\nTarget: Público general en Yangon\nTáctica: Descuentos cruzados con Sports & Travel para liquidar stock en marzo.", 9, 3, 3.5, 2, 18)

    add_insight_box(slide, "Data al servicio de la acción: no le hablamos a todos por igual, segmentamos la oferta.")

    prs.save('/Users/facundo/Desktop/tallerStoryteling/Taller/Archivos notables/Desafio3_Presentacion_Redisenada.pptx')
    print("Presentación creada exitosamente.")

if __name__ == '__main__':
    create_modern_ppt()
