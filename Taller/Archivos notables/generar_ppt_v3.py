from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import pptx.enum.text as pptxt

SW, SH = 13.333, 7.5

# Paleta: verde bosque + terracota + crema
BG      = RGBColor(250, 249, 246)   # crema cálido
GREEN   = RGBColor(30, 60, 48)      # verde bosque profundo
SAGE    = RGBColor(100, 140, 110)   # verde salvia
TERRA   = RGBColor(178, 90, 55)     # terracota cálida
CREAM   = RGBColor(242, 238, 228)   # crema suave
TEXT    = RGBColor(25, 25, 22)      # casi negro
MUTED   = RGBColor(130, 128, 118)   # gris cálido
PH_BG   = RGBColor(234, 238, 232)   # placeholder bg
PH_BD   = RGBColor(170, 185, 165)   # placeholder borde
WHITE   = RGBColor(255, 255, 255)
LGREEN  = RGBColor(220, 232, 222)   # verde claro para fondos

def prs():
    p = Presentation()
    p.slide_width  = Inches(SW)
    p.slide_height = Inches(SH)
    return p

def blank(pr):
    return pr.slides.add_slide(pr.slide_layouts[6])

def bg(slide, c=BG):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = c

def R(slide, l, t, w, h, fc, lc=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fc
    if lc:
        s.line.color.rgb = lc; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def T(slide, text, l, t, w, h, sz=16, bold=False, color=TEXT,
      align=PP_ALIGN.LEFT, italic=False, wrap=True):
    b = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.bold = bold
    p.font.italic = italic; p.font.color.rgb = color
    p.font.name = 'Calibri'; p.alignment = align
    return b

def PH(slide, l, t, w, h, label="Insertar gráfico aquí"):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = PH_BG
    s.line.color.rgb = PH_BD; s.line.width = Pt(1.2)
    tf = s.text_frame
    tf.vertical_anchor = pptxt.MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label; p.font.size = Pt(12); p.font.italic = True
    p.font.color.rgb = MUTED; p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

def footer(slide):
    T(slide, "Desafío 3  ·  Supermarket Sales  ·  Myanmar 2019",
      0.5, 7.2, 12, 0.28, sz=9, color=MUTED, italic=True)

def slide_header(slide, title, subtitle, accent_color=TERRA):
    R(slide, 0, 0, SW, 0.07, accent_color)
    T(slide, title, 0.6, 0.3, 11, 0.7, sz=32, bold=True, color=GREEN)
    T(slide, subtitle, 0.6, 0.98, 11, 0.38, sz=15, italic=True, color=MUTED)

pr = prs()

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ═══════════════════════════════════════════════════════════════
s = blank(pr); bg(s, GREEN)
R(s, 0, 0, SW, 0.08, TERRA)
# Bloque crema en mitad derecha para romper el bloque sólido
R(s, 8.5, 0, 4.83, SH, RGBColor(38, 68, 55))

T(s, "MYANMAR · 2019", 0.7, 0.8, 8, 0.42, sz=11, bold=True,
  color=TERRA, italic=False)
T(s, "Supermarket\nSales Analysis", 0.7, 1.3, 9, 2.8,
  sz=64, bold=True, color=WHITE)
T(s, "Análisis de datos para 3 campañas de marketing mensuales",
  0.7, 4.05, 7.5, 0.7, sz=19, italic=True,
  color=RGBColor(170, 200, 175))
R(s, 0.7, 4.85, 4.2, 0.055, TERRA)
T(s, "Cliente: Cadena de Supermercados  |  3 sucursales  |  Enero – Marzo",
  0.7, 4.95, 9, 0.4, sz=11, color=RGBColor(140, 170, 145), italic=True)

# KPIs
for i, (val, lbl) in enumerate([("1,000","Transacciones"),
                                  ("$322.9K","Ventas totales"),
                                  ("$323","Ticket promedio")]):
    x = 9.0 + (i % 3) * 0.01
    y = 1.2 + i * 1.7
    R(s, 8.8, y, 4.1, 1.45, RGBColor(22, 50, 38))
    T(s, val, 8.8, y + 0.08, 4.1, 0.85, sz=42, bold=True,
      color=TERRA, align=PP_ALIGN.CENTER)
    T(s, lbl, 8.8, y + 0.92, 4.1, 0.4, sz=13,
      color=RGBColor(160, 195, 165), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — CONTEXTO
# ═══════════════════════════════════════════════════════════════
s = blank(pr); bg(s, BG)
slide_header(s, "Contexto del Negocio",
             "Tres ciudades, tres perfiles de consumidor en Myanmar")

cities = [
    ("A", "Yangon",    "Capital comercial y económica", "5.9M hab.",
     "340 transacciones · Ticket prom: $312 · Rating 7.0",
     "La ciudad más cosmopolita. Mayor volumen pero ticket más bajo. El consumidor es el más diverso y competido.", GREEN),
    ("B", "Mandalay",  "Centro cultural del norte",     "1.6M hab.",
     "332 transacciones · Ticket prom: $320 · Rating 6.8",
     "Hub religioso y cultural. Lidera en Health & Beauty y Sports & Travel. Menor satisfacción del cliente.", SAGE),
    ("C", "Naypyitaw", "Capital política y gubernamental","815K hab.",
     "328 transacciones · Ticket prom: $337 · Rating 7.1",
     "La ciudad más chica, el ticket más alto. Los funcionarios de gobierno tienen ingresos estables y compran más por visita.", TERRA),
]

for i, (letter, city, role, pop, stats, desc, col) in enumerate(cities):
    x = 0.5 + i * 4.25; cw = 4.0
    R(s, x, 1.65, cw, 5.2, WHITE, lc=LGREEN, lw=0.8)
    R(s, x, 1.65, 0.28, 5.2, col)
    T(s, letter, x + 0.38, 1.75, 0.6, 0.55, sz=22, bold=True, color=col)
    T(s, city,   x + 0.38, 2.35, cw - 0.55, 0.48, sz=20, bold=True, color=GREEN)
    T(s, role,   x + 0.38, 2.83, cw - 0.55, 0.38, sz=12, bold=True, color=TERRA)
    T(s, pop,    x + 0.38, 3.22, cw - 0.55, 0.32, sz=11, italic=True, color=MUTED)
    R(s, x + 0.38, 3.6, cw - 0.55, 0.04, LGREEN)
    T(s, stats,  x + 0.38, 3.7, cw - 0.55, 0.45, sz=11, bold=True, color=GREEN)
    T(s, desc,   x + 0.38, 4.2, cw - 0.55, 2.3,  sz=12, color=TEXT, wrap=True)

footer(s)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — GÉNERO  (layout: dos stats grandes arriba, gráfico abajo)
# ═══════════════════════════════════════════════════════════════
s = blank(pr); bg(s, BG)
slide_header(s, "Comportamiento por Género",
             "¿Cómo se comparan las compras de mujeres vs hombres?", SAGE)

# Dos bloques de stat lado a lado
for i, (val, sub, detail, col) in enumerate([
    ("$335", "ticket prom. · Mujeres",
     "501 transacciones · $167.8K totales · 52% de las ventas", TERRA),
    ("$311", "ticket prom. · Hombres",
     "499 transacciones · $155.1K totales · 48% de las ventas", SAGE),
]):
    x = 0.6 + i * 6.4; bw = 5.9
    R(s, x, 1.58, bw, 2.5, CREAM, lc=LGREEN, lw=0.8)
    R(s, x, 1.58, bw, 0.06, col)
    T(s, val, x + 0.25, 1.72, bw - 0.4, 1.15,
      sz=66, bold=True, color=col)
    T(s, sub, x + 0.25, 2.78, bw - 0.4, 0.38,
      sz=14, bold=True, color=GREEN)
    T(s, detail, x + 0.25, 3.15, bw - 0.4, 0.38,
      sz=11, italic=True, color=MUTED)

# Insight strip
R(s, 0.6, 4.25, 12.1, 0.72, RGBColor(30, 60, 48))
T(s, "Insight: En Health & Beauty, los hombres gastan un 65% más que las mujeres — la mayor brecha de todo el dataset.",
  0.9, 4.38, 11.5, 0.48, sz=15, bold=True, color=WHITE)

PH(s, 0.6, 5.1, 12.1, 1.75, "Insertar gráfico: ventas y ticket promedio por género")
footer(s)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — HORARIOS  (layout: gráfico izq grande, insight panel der)
# ═══════════════════════════════════════════════════════════════
s = blank(pr); bg(s, BG)
slide_header(s, "El Ritmo del Día",
             "¿En qué momentos se concentran las ventas?")

PH(s, 0.6, 1.58, 7.8, 5.25, "Insertar gráfico: ventas y transacciones por hora (10h – 20h)")

# Panel derecho
R(s, 8.6, 1.58, 4.1, 5.25, GREEN)
R(s, 8.6, 1.58, 4.1, 0.06, TERRA)
T(s, "DOS PICOS CLAVE", 8.85, 1.72, 3.65, 0.38,
  sz=10, bold=True, color=TERRA)
T(s, "13:00 hs", 8.85, 2.18, 3.65, 0.6, sz=34, bold=True, color=WHITE)
T(s, "$34.7K · 103 transacciones", 8.85, 2.75, 3.65, 0.38,
  sz=12, color=RGBColor(180, 210, 185))
R(s, 8.85, 3.25, 3.3, 0.04, RGBColor(80, 115, 90))
T(s, "19:00 hs", 8.85, 3.4, 3.65, 0.6, sz=34, bold=True, color=TERRA)
T(s, "$39.7K · 113 transacciones", 8.85, 3.97, 3.65, 0.38,
  sz=12, color=RGBColor(180, 210, 185))
R(s, 8.85, 4.48, 3.3, 0.04, RGBColor(80, 115, 90))
T(s, "Valle: 16h, 17h y 20h están por debajo del promedio. Oportunidad para generar tráfico con ofertas.",
  8.85, 4.62, 3.55, 1.8, sz=12, color=RGBColor(190, 215, 192), wrap=True)
footer(s)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — LÍNEA DE PRODUCTOS  (layout: gráfico izq, ranking + insight der)
# ═══════════════════════════════════════════════════════════════
s = blank(pr); bg(s, BG)
slide_header(s, "Desempeño del Portafolio",
             "¿Qué categorías generan más ingresos?", TERRA)

PH(s, 0.6, 1.58, 7.4, 5.25,
   "Insertar gráfico: ingresos por línea de producto (barras horizontales)")

# Ranking lateral
R(s, 8.2, 1.58, 4.5, 5.25, CREAM, lc=LGREEN, lw=0.8)
T(s, "RANKING DE INGRESOS", 8.4, 1.72, 4.1, 0.35,
  sz=10, bold=True, color=TERRA)
ranking = [
    ("1", "Food & Beverages",     "$56.1K", True),
    ("2", "Sports & Travel",       "$55.1K", False),
    ("3", "Electronic Acc.",       "$54.3K", False),
    ("4", "Fashion Acc.",          "$54.3K", False),
    ("5", "Home & Lifestyle",      "$53.9K", False),
    ("6", "Health & Beauty",       "$49.1K", True),
]
for j, (num, cat, amt, highlight) in enumerate(ranking):
    y = 2.15 + j * 0.72
    fc = RGBColor(242, 228, 218) if highlight else WHITE
    lc2 = TERRA if highlight else LGREEN
    R(s, 8.35, y, 4.1, 0.62, fc, lc=lc2, lw=0.8)
    T(s, num, 8.45, y + 0.08, 0.35, 0.44, sz=14, bold=True, color=MUTED)
    T(s, cat, 8.82, y + 0.08, 2.4, 0.44, sz=13,
      bold=highlight, color=GREEN)
    T(s, amt, 11.2, y + 0.08, 1.15, 0.44, sz=13, bold=True,
      color=TERRA if highlight else TEXT, align=PP_ALIGN.RIGHT)
footer(s)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — PRODUCTO × GÉNERO  (layout: gráfico top, brechas abajo)
# ═══════════════════════════════════════════════════════════════
s = blank(pr); bg(s, BG)
slide_header(s, "Productos por Género",
             "¿Qué compra cada género por categoría?", SAGE)

PH(s, 0.6, 1.58, 12.1, 3.1,
   "Insertar gráfico: % por género en cada línea de producto (barras 100% apiladas)")

# Mini stats de brechas en fila inferior
brechas = [
    ("Health & Beauty",     "Hombres 62.3%",  "Mujeres 37.7%",  TERRA, SAGE),
    ("Food & Beverages",    "Mujeres 59.1%",  "Hombres 40.9%",  SAGE,  TERRA),
    ("Fashion Acc.",        "Mujeres 56.0%",  "Hombres 44.0%",  SAGE,  TERRA),
    ("Electronic Acc.",     "Casi igualado",  "50% / 50%",       MUTED, MUTED),
]
for i, (cat, v1, v2, c1, c2) in enumerate(brechas):
    x = 0.6 + i * 3.1; bw = 2.85
    R(s, x, 4.88, bw, 1.95, WHITE, lc=LGREEN, lw=0.8)
    R(s, x, 4.88, bw, 0.06, c1)
    T(s, cat, x + 0.15, 4.98, bw - 0.25, 0.38, sz=12, bold=True, color=GREEN)
    T(s, v1,  x + 0.15, 5.38, bw - 0.25, 0.38, sz=13, bold=True, color=c1)
    T(s, v2,  x + 0.15, 5.75, bw - 0.25, 0.38, sz=13, bold=True, color=c2)

R(s, 12.7, 4.88, 0.6, 1.95, GREEN)   # decorative right edge block
footer(s)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — SUCURSALES  (layout: 3 cards arriba, gráfico abajo)
# ═══════════════════════════════════════════════════════════════
s = blank(pr); bg(s, BG)
slide_header(s, "Performance por Sucursal",
             "¿Cómo se desempeña cada sucursal?")

branch_data = [
    ("A · Yangon",    "$106.2K", "340 transacciones · ticket $312 · rating 7.0",
     "Mayor volumen, menor ticket. Lidera en Home & Lifestyle.", GREEN),
    ("C · Naypyitaw", "$110.6K", "328 transacciones · ticket $337 · rating 7.1",
     "Menor volumen pero mayor ticket y mejor satisfacción.", TERRA),
    ("B · Mandalay",  "$106.2K", "332 transacciones · ticket $320 · rating 6.8",
     "Lidera en Health & Beauty. Rating más bajo: oportunidad de mejora.", SAGE),
]
for i, (name, total, stats, desc, col) in enumerate(branch_data):
    x = 0.6 + i * 4.25; bw = 4.0
    R(s, x, 1.55, bw, 2.85, WHITE, lc=LGREEN, lw=0.8)
    R(s, x, 1.55, bw, 0.07, col)
    T(s, name,  x + 0.2, 1.68, bw - 0.35, 0.42, sz=14, bold=True, color=GREEN)
    T(s, total, x + 0.2, 2.1,  bw - 0.35, 0.82, sz=38, bold=True, color=col)
    T(s, stats, x + 0.2, 2.9,  bw - 0.35, 0.38, sz=11, italic=True, color=MUTED)
    T(s, desc,  x + 0.2, 3.3,  bw - 0.35, 0.9,  sz=12, color=TEXT, wrap=True)

PH(s, 0.6, 4.62, 12.1, 2.2,
   "Insertar gráfico: ventas y transacciones por sucursal")
footer(s)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — CAMPAÑAS  (3 columnas iguales, texto controlado)
# ═══════════════════════════════════════════════════════════════
s = blank(pr); bg(s, BG)
R(s, 0, 0, SW, 0.07, TERRA)
T(s, "Propuesta de Campañas de Marketing", 0.6, 0.25, 11, 0.72,
  sz=32, bold=True, color=GREEN)
T(s, "Una acción concreta por mes, fundamentada en los datos del dataset",
  0.6, 0.95, 11, 0.42, sz=15, italic=True, color=MUTED)

campaigns = [
    ("MES 1 · ENERO", "Sabores y Familia", GREEN, [
        ("Foco",   "Food & Beverages"),
        ("Target", "Mujeres — Naypyitaw y Yangon"),
        ("Táctica","Degustaciones y combos familiares en los picos de 13h y 19h."),
        ("Por qué","Enero es el mejor mes ($116K). Mujeres dominan esta categoría (+44%). Naypyitaw lidera el mes."),
    ]),
    ("MES 2 · FEBRERO", "Bienestar para Él", TERRA, [
        ("Foco",   "Health & Beauty"),
        ("Target", "Hombres — Mandalay"),
        ("Táctica","Comunicación y promociones pensadas exclusivamente para hombres."),
        ("Por qué","Febrero es el mes más bajo ($97K). La brecha masculina en H&B es la mayor del dataset: 65%."),
    ]),
    ("MES 3 · MARZO", "Hogar y Estilo", SAGE, [
        ("Foco",   "Home & Lifestyle + Sports & Travel"),
        ("Target", "Público general — Yangon"),
        ("Táctica","Combos cruzados hogar + deporte. Aprovechar el ticket alto de Yangon en estas líneas."),
        ("Por qué","Marzo antecede al Año Nuevo birmano (abril): las familias renuevan el hogar."),
    ]),
]

for i, (mes, title, col, items) in enumerate(campaigns):
    x = 0.5 + i * 4.25; cw = 4.0
    # Card fondo
    R(s, x, 1.58, cw, 5.7, WHITE, lc=LGREEN, lw=0.8)
    # Header band
    R(s, x, 1.58, cw, 1.0, col)
    T(s, mes,   x + 0.2, 1.65, cw - 0.3, 0.35, sz=9,  bold=True,
      color=WHITE if col != SAGE else RGBColor(25, 50, 35))
    T(s, title, x + 0.2, 2.0,  cw - 0.3, 0.5,  sz=18, bold=True,
      color=WHITE if col != SAGE else RGBColor(25, 50, 35))
    # Items
    for j, (label, detail) in enumerate(items):
        y = 2.8 + j * 1.1
        T(s, label.upper(), x + 0.2, y,      cw - 0.3, 0.32, sz=9,
          bold=True, color=col)
        T(s, detail,        x + 0.2, y + 0.32, cw - 0.3, 0.72, sz=11,
          color=TEXT, wrap=True)

footer(s)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — CIERRE
# ═══════════════════════════════════════════════════════════════
s = blank(pr); bg(s, GREEN)
R(s, 0, 0, SW, 0.07, TERRA)
R(s, 9.5, 0, 3.83, SH, RGBColor(22, 50, 38))   # bloque derecho sutil

T(s, "LO QUE NOS DICE LA DATA", 0.7, 0.8, 8.5, 0.45,
  sz=11, bold=True, color=TERRA, italic=False)
T(s, "Conclusiones", 0.7, 1.3, 8.5, 1.0,
  sz=52, bold=True, color=WHITE)
R(s, 0.7, 2.45, 4.5, 0.06, TERRA)

pts = [
    "La distribución entre géneros, sucursales y horarios es equilibrada — pero con brechas muy accionables.",
    "Health & Beauty: los hombres gastan 65% más que las mujeres. El mayor insight del dataset.",
    "Los picos de 13h y 19h deben ser el eje de cualquier activación en sucursal.",
    "Naypyitaw factura más con menos transacciones: prioridad para campañas de ticket alto.",
    "Febrero es el mes más débil. Una campaña focalizada (Health & Beauty / hombres) puede revertirlo.",
]
for i, pt in enumerate(pts):
    y = 2.65 + i * 0.82
    R(s, 0.7, y + 0.1, 0.22, 0.22, TERRA)
    T(s, pt, 1.1, y, 8.1, 0.75, sz=15, color=WHITE, wrap=True)

T(s, "Desafío 3  ·  Supermarket Sales  ·  Myanmar 2019",
  0.7, 7.15, 12, 0.3, sz=9, color=RGBColor(100, 135, 108), italic=True)

# ═══════════════════════════════════════════════════════════════
out = ('/Users/facundo/Desktop/tallerStoryteling/Taller/'
       'Archivos notables/Desafio3_Presentacion_v3.pptx')
pr.save(out)
print(f"OK: {out}")
