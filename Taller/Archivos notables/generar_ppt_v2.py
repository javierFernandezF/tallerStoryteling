from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import pptx.enum.text as pptxt

# === DIMENSIONES ===
SW = 13.333
SH = 7.5
SBW = 4.8  # sidebar width

# === PALETA ===
NAVY  = RGBColor(18, 44, 77)
AMBER = RGBColor(212, 148, 26)
TEAL  = RGBColor(30, 140, 140)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(245, 247, 250)
TEXT  = RGBColor(26, 26, 46)
MUTED = RGBColor(110, 122, 140)
PH_BG = RGBColor(238, 243, 250)
PH_BD = RGBColor(170, 195, 225)

def new_prs():
    p = Presentation()
    p.slide_width = Inches(SW)
    p.slide_height = Inches(SH)
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fc, lc=None, lw=0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fc
    if lc:
        s.line.color.rgb = lc
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def add_txt(slide, text, l, t, w, h, size=16, bold=False, color=TEXT,
            align=PP_ALIGN.LEFT, italic=False, wrap=True):
    b = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.alignment = align
    return b

def add_multiline_txt(slide, lines, l, t, w, h, size=15, color=WHITE, bold_first=False):
    b = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.font.bold = (i == 0 and bold_first)
        if line == "":
            p.font.size = Pt(6)

def chart_ph(slide, l, t, w, h, label):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = PH_BG
    s.line.color.rgb = PH_BD
    s.line.width = Pt(1.5)
    tf = s.text_frame
    tf.vertical_anchor = pptxt.MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = MUTED
    p.font.italic = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

def footer(slide):
    add_rect(slide, 0, 7.15, SW, 0.35, LIGHT)
    add_txt(slide, "Desafío 3  ·  Supermarket Sales  ·  Myanmar 2019",
            0.4, 7.18, 10, 0.28, size=9, color=MUTED, italic=True)

def data_slide(prs, tag, kpi, title, question, insight_lines, ph_label):
    slide = blank(prs)
    set_bg(slide, WHITE)

    # Sidebar navy
    add_rect(slide, 0, 0, SBW, SH, NAVY)
    # Amber left bar
    add_rect(slide, 0, 0, 0.2, SH, AMBER)
    # Amber top bar (right side only)
    add_rect(slide, SBW, 0, SW - SBW, 0.08, AMBER)

    # Tag label
    add_txt(slide, tag.upper(), 0.35, 0.45, SBW - 0.5, 0.38,
            size=10, bold=True, color=AMBER)
    # KPI big number
    add_txt(slide, kpi, 0.35, 0.9, SBW - 0.4, 1.35,
            size=60, bold=True, color=AMBER)
    # Divider
    add_rect(slide, 0.35, 2.4, 3.5, 0.05, AMBER)
    # Title
    add_txt(slide, title, 0.35, 2.55, SBW - 0.45, 1.5,
            size=22, bold=True, color=WHITE, wrap=True)
    # Insight
    add_multiline_txt(slide, insight_lines, 0.35, 4.2, SBW - 0.45, 2.5,
                      size=14, color=WHITE)

    # Right side
    add_txt(slide, question, SBW + 0.35, 0.22, SW - SBW - 0.6, 0.5,
            size=12, italic=True, color=MUTED)
    chart_ph(slide, SBW + 0.35, 0.85, SW - SBW - 0.6, 5.9, ph_label)
    footer(slide)
    return slide


# =========================================================
# CONSTRUIR LA PRESENTACIÓN
# =========================================================
prs = new_prs()

# --- SLIDE 1: PORTADA ---
slide = blank(prs)
set_bg(slide, NAVY)
add_rect(slide, 0, 0, 0.22, SH, AMBER)

add_txt(slide, "MYANMAR · 2019", 0.55, 0.7, 11, 0.45,
        size=12, color=AMBER, italic=True, bold=True)
add_txt(slide, "Supermarket Sales", 0.55, 1.2, 12, 1.8,
        size=62, bold=True, color=WHITE)
add_txt(slide, "Análisis de datos para 3 campañas de marketing mensuales",
        0.55, 2.85, 11, 0.65, size=22, color=RGBColor(170, 195, 230), italic=True)
add_rect(slide, 0.55, 3.55, 5.0, 0.07, AMBER)

# KPIs
kpis = [("1,000", "Transacciones"), ("$322.9K", "Ventas totales"), ("$323", "Ticket promedio")]
for i, (val, lbl) in enumerate(kpis):
    x = 0.55 + i * 4.1
    add_rect(slide, x, 4.1, 3.6, 2.2, RGBColor(25, 55, 95))
    add_rect(slide, x, 4.1, 3.6, 0.07, AMBER)
    add_txt(slide, val, x, 4.3, 3.6, 1.1, size=48, bold=True,
            color=AMBER, align=PP_ALIGN.CENTER)
    add_txt(slide, lbl, x, 5.35, 3.6, 0.45, size=15,
            color=RGBColor(170, 195, 230), align=PP_ALIGN.CENTER)

add_txt(slide, "Cliente: Cadena de Supermercados  |  3 sucursales  |  Enero – Marzo 2019",
        0.55, 6.95, 12, 0.4, size=10, color=MUTED, italic=True)

# --- SLIDE 2: CONTEXTO ---
slide = blank(prs)
set_bg(slide, LIGHT)
add_rect(slide, 0, 0, 0.22, SH, AMBER)
add_rect(slide, 0, 0, SW, 0.08, AMBER)

add_txt(slide, "Contexto del Negocio", 0.55, 0.4, 10, 0.75,
        size=34, bold=True, color=NAVY)
add_txt(slide, "Tres ciudades con perfiles de consumidor muy distintos",
        0.55, 1.1, 10, 0.45, size=17, color=MUTED, italic=True)

cities = [
    ("A", "Yangon",     "Capital comercial",   "5.9M hab.",
     ["340 transacciones", "Ticket prom: $312", "Rating: 7.0", "",
      "La ciudad más cosmopolita del país. Mayor volumen de transacciones pero ticket promedio más bajo. El consumidor es el más diverso."]),
    ("B", "Mandalay",   "Centro cultural",     "1.6M hab.",
     ["332 transacciones", "Ticket prom: $320", "Rating: 6.8", "",
      "Hub religioso y cultural. Lidera en Health & Beauty y Sports & Travel. Tiene el menor índice de satisfacción."]),
    ("C", "Naypyitaw",  "Capital política",    "815K hab.",
     ["328 transacciones", "Ticket prom: $337", "Rating: 7.1", "",
      "La ciudad más pequeña pero la de mayor ticket. Los funcionarios del gobierno tienen ingresos estables."]),
]

for i, (letter, city, role, pop, details) in enumerate(cities):
    x = 0.55 + i * 4.25
    cw = 4.0
    add_rect(slide, x, 1.85, cw, 5.0, WHITE, lc=PH_BD, lw=0.5)
    add_rect(slide, x, 1.85, cw, 0.75, NAVY)
    add_rect(slide, x, 1.85, 0.7, 0.75, AMBER)
    add_txt(slide, letter, x, 1.88, 0.7, 0.65, size=22, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
    add_txt(slide, city, x + 0.75, 1.9, cw - 0.85, 0.4,
            size=20, bold=True, color=WHITE)
    add_txt(slide, role, x + 0.15, 2.65, cw - 0.3, 0.35,
            size=12, bold=True, color=AMBER)
    add_txt(slide, pop, x + 0.15, 3.0, cw - 0.3, 0.35,
            size=11, italic=True, color=MUTED)
    add_multiline_txt(slide, details, x + 0.15, 3.35, cw - 0.3, 3.0,
                      size=12, color=TEXT, bold_first=True)

footer(slide)

# --- SLIDE 3: GÉNERO ---
data_slide(prs,
    tag="Análisis de Género",
    kpi="8% más",
    title="Las mujeres gastan más por visita. Los hombres dominan un nicho clave.",
    question="¿Cómo se comparan las compras de mujeres vs hombres?",
    insight_lines=[
        "Ticket promedio:", "Mujeres: $335 · Hombres: $311", "",
        "Las mujeres generan el 52% de las ventas ($167.8K vs $155K).", "",
        "El dato más relevante: en Health & Beauty, los hombres gastan un 65% más que las mujeres.",
        "Una oportunidad casi sin explotar."
    ],
    ph_label="Insertar gráfico: Ventas totales y ticket promedio por género"
)

# --- SLIDE 4: HORARIOS ---
data_slide(prs,
    tag="Distribución Horaria",
    kpi="19:00 hs",
    title="El pico de la tarde-noche domina. Dos momentos concentran el tráfico.",
    question="¿En qué momentos del día se concentran las ventas?",
    insight_lines=[
        "Pico 1 — 13:00 hs: $34.7K, 103 transacciones.", "",
        "Pico 2 — 19:00 hs: $39.7K, 113 transacciones.", "(momento de mayor volumen absoluto)", "",
        "Horas valle: 16h, 17h y 20h están por debajo del promedio horario ($29.4K).",
        "Ahí hay margen para generar tráfico con ofertas dirigidas."
    ],
    ph_label="Insertar gráfico: Ventas y transacciones por hora (10h a 20h)"
)

# --- SLIDE 5: LÍNEA DE PRODUCTOS ---
data_slide(prs,
    tag="Portafolio de Productos",
    kpi="$56.1K",
    title="Los ingresos están distribuidos. Food & Beverages lidera.",
    question="¿Qué categorías generan más ingresos?",
    insight_lines=[
        "Food & Beverages:     $56.1K  (1°)",
        "Sports & Travel:        $55.1K  (2°)",
        "Electronic Acc.:         $54.3K  (3°)",
        "Fashion Acc.:             $54.3K  (4°)",
        "Home & Lifestyle:     $53.9K  (5°)", "",
        "Health & Beauty:      $49.1K  (último)", "",
        "Fashion tiene la mayor cant. de transacciones (178) pero ticket más bajo."
    ],
    ph_label="Insertar gráfico: Ingresos por línea de producto (barras horizontales)"
)

# --- SLIDE 6: PRODUCTO × GÉNERO ---
data_slide(prs,
    tag="Producto por Género",
    kpi="65%",
    title="Cada categoría tiene un perfil de cliente distinto. El género importa.",
    question="¿Qué compra cada género por línea de producto?",
    insight_lines=[
        "Health & Beauty → Hombres 62.3% / Mujeres 37.7%",
        "Food & Beverages → Mujeres 59.1% / Hombres 40.9%",
        "Fashion Acc. → Mujeres 56% / Hombres 44%",
        "Home & Lifestyle → Mujeres 55.8% / Hombres 44.2%",
        "Sports & Travel → Casi igualado (52% / 48%)",
        "Electronic Acc. → Igualado (50% / 50%)", "",
        "La brecha más grande: Health & Beauty (24.5 puntos)."
    ],
    ph_label="Insertar gráfico: % por género en cada línea de producto (barras 100% apiladas)"
)

# --- SLIDE 7: SUCURSALES ---
data_slide(prs,
    tag="Performance por Sucursal",
    kpi="Suc. C",
    title="Naypyitaw lidera en ventas y satisfacción. Yangon tiene más volumen.",
    question="¿Cómo se desempeña cada sucursal?",
    insight_lines=[
        "C (Naypyitaw): $110.6K · ticket $337 · rating 7.1",
        "A (Yangon):      $106.2K · ticket $312 · rating 7.0",
        "B (Mandalay):  $106.2K · ticket $320 · rating 6.8", "",
        "Naypyitaw factura más con menos transacciones.",
        "Yangon tiene el mayor volumen pero el menor ticket.",
        "Mandalay tiene el menor índice de satisfacción: oportunidad de mejora."
    ],
    ph_label="Insertar gráfico: Ventas y transacciones por sucursal"
)

# --- SLIDE 8: CAMPAÑAS ---
slide = blank(prs)
set_bg(slide, LIGHT)
add_rect(slide, 0, 0, 0.22, SH, AMBER)
add_rect(slide, 0, 0, SW, 0.08, AMBER)

add_txt(slide, "Propuesta de Campañas de Marketing", 0.55, 0.35, 11, 0.75,
        size=34, bold=True, color=NAVY)
add_txt(slide, "Una acción por mes, fundamentada en los datos",
        0.55, 1.05, 10, 0.45, size=17, color=MUTED, italic=True)

cdata = [
    ("MES 1 · ENERO", "Sabores\ny Familia",   NAVY,
     ["Foco: Food & Beverages",
      "Target: Mujeres, Naypyitaw y Yangon",
      "Táctica: Degustaciones y combos familiares en los picos de 13h y 19h.",
      "",
      "Enero es el mejor mes ($116K). Las mujeres dominan esta categoría (+44%). Naypyitaw lidera en enero."]),
    ("MES 2 · FEBRERO", "Bienestar\npara Él", TEAL,
     ["Foco: Health & Beauty",
      "Target: Hombres, Mandalay",
      "Táctica: Comunicación y descuentos pensados para hombres. Reforzar la sucursal B.",
      "",
      "Febrero es el mes más bajo ($97K). La brecha masculina en Health & Beauty es la mayor del dataset: 65%."]),
    ("MES 3 · MARZO", "Hogar\ny Estilo",      AMBER,
     ["Foco: Home & Lifestyle + Sports & Travel",
      "Target: Público general, Yangon",
      "Táctica: Combos cruzados hogar + deporte. Yangon domina esta categoría.",
      "",
      "Marzo se recupera ($109K) y antecede al Año Nuevo birmano (abril), cuando las familias renuevan el hogar."]),
]

for i, (mes, title, hcolor, details) in enumerate(cdata):
    x = 0.55 + i * 4.25
    cw = 4.0
    add_rect(slide, x, 1.75, cw, 5.0, WHITE, lc=PH_BD, lw=0.5)
    add_rect(slide, x, 1.75, cw, 0.9, hcolor)
    add_txt(slide, mes, x + 0.15, 1.8, cw - 0.3, 0.38,
            size=10, bold=True, color=WHITE if hcolor != AMBER else NAVY)
    add_txt(slide, title, x + 0.15, 2.15, cw - 0.3, 0.65,
            size=20, bold=True, color=WHITE if hcolor != AMBER else NAVY)
    add_multiline_txt(slide, details, x + 0.15, 2.75, cw - 0.3, 3.7,
                      size=12, color=TEXT, bold_first=True)

footer(slide)

# --- SLIDE 9: CIERRE ---
slide = blank(prs)
set_bg(slide, NAVY)
add_rect(slide, 0, 0, 0.22, SH, AMBER)
add_rect(slide, 0, 0, SW, 0.08, AMBER)

add_txt(slide, "CONCLUSIONES", 0.55, 0.7, 11, 0.5,
        size=13, bold=True, color=AMBER, italic=False)
add_txt(slide, "Los datos hablan claro.", 0.55, 1.25, 11, 1.0,
        size=48, bold=True, color=WHITE)
add_rect(slide, 0.55, 2.35, 5.0, 0.07, AMBER)

conclusions = [
    "La distribución entre géneros y sucursales es equilibrada, pero con brechas accionables.",
    "Health & Beauty es el producto con mayor oportunidad: los hombres gastan 65% más que las mujeres.",
    "Los picos de 13h y 19h deben ser el eje de cualquier activación en sucursal.",
    "Naypyitaw es la sucursal más eficiente: factura más con menos transacciones.",
    "Febrero es el mes más débil. Una campaña focalizada puede revertirlo.",
]

for i, c in enumerate(conclusions):
    y = 2.6 + i * 0.75
    add_rect(slide, 0.55, y + 0.12, 0.25, 0.25, AMBER)
    add_txt(slide, c, 0.95, y, 11.5, 0.65, size=16, color=WHITE)

add_txt(slide, "Desafío 3  ·  Supermarket Sales  ·  Myanmar 2019",
        0.55, 7.1, 12, 0.35, size=9, color=MUTED, italic=True)

# --- GUARDAR ---
out = '/Users/facundo/Desktop/tallerStoryteling/Taller/Archivos notables/Desafio3_Presentacion_v2.pptx'
prs.save(out)
print(f"Guardado en: {out}")
