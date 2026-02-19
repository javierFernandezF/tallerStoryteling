from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd

# Cargar datos
def load_data():
    data_rows = []
    with open('/Users/javierfernandez/Desktop/Taller/Libro1.csv', 'r', encoding='utf-8-sig') as f:
        lines = f.readlines()
    header_line = lines[0].strip()
    if header_line.startswith('"') and header_line.endswith('"'):
        header = header_line[1:-1].split(';')
    for line in lines[1:]:
        line = line.strip()
        if line.startswith('"') and line.endswith('"'):
            row_data = line[1:-1].split(';')
            data_rows.append(row_data)
    df = pd.DataFrame(data_rows, columns=header)
    numeric_cols = ['Unit price', 'Quantity', 'Tax 5%', 'Total', 'Time', 'cogs',
                    'gross margin percentage', 'gross income', 'Rating']
    for col in numeric_cols:
        if col in df.columns:
            df[col] = df[col].str.replace(',', '.').astype(float)
    df['Hour'] = (df['Time'] * 24).round().astype(int).clip(0, 23)
    return df

df = load_data()

# Paleta de colores
AZUL_OSCURO  = RGBColor(0x1A, 0x37, 0x5E)
AZUL_MEDIO   = RGBColor(0x27, 0x6F, 0xBF)
VERDE_ACENTO = RGBColor(0x2E, 0xCC, 0x71)
ROJO_ACENTO  = RGBColor(0xE7, 0x4C, 0x3C)
GRIS_CLARO   = RGBColor(0xF4, 0xF6, 0xF9)
GRIS_TEXTO   = RGBColor(0x55, 0x55, 0x55)
BLANCO       = RGBColor(0xFF, 0xFF, 0xFF)

# Helpers
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h, font_size=12, bold=False,
                color=None, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox

def add_placeholder_box(slide, l, t, w, h, label="[ Insertar gráfico de Excel aquí ]"):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEC, 0xF0, 0xF1)
    shape.line.color.rgb = AZUL_MEDIO
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(13)
    run.font.color.rgb = AZUL_MEDIO
    run.font.italic = True
    return shape

def kpi_box(slide, l, t, w, h, value, label, color):
    add_rect(slide, l, t, w, h, color)
    add_textbox(slide, value, l, t + 0.08, w, h * 0.52,
                font_size=22, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, l, t + h * 0.55, w, h * 0.42,
                font_size=10, color=BLANCO, align=PP_ALIGN.CENTER)

def slide_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 10, 1.2, AZUL_OSCURO)
    add_textbox(slide, title, 0.3, 0.15, 9.4, 0.7,
                font_size=24, bold=True, color=BLANCO)
    if subtitle:
        add_textbox(slide, subtitle, 0.3, 0.78, 9.4, 0.35,
                    font_size=12, color=RGBColor(0xAA, 0xCC, 0xFF))

def slide_footer(slide):
    add_rect(slide, 0, 7.28, 10, 0.22, AZUL_OSCURO)
    add_textbox(slide, "Desafío 3  ·  Supermarket Sales Analysis  ·  2026",
                0.2, 7.28, 9.6, 0.22, font_size=9,
                color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.RIGHT)

# Crear presentación nueva
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# ============================================================
# SLIDE 1 — PORTADA
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, AZUL_OSCURO)
add_rect(slide, 0, 0, 0.35, 7.5, AZUL_MEDIO)

add_textbox(slide, "Supermarket Sales", 0.7, 1.5, 8.8, 1.1,
            font_size=40, bold=True, color=BLANCO)
add_textbox(slide, "Análisis de Datos para Campañas de Marketing",
            0.7, 2.6, 8.8, 0.6, font_size=18, color=RGBColor(0xAA, 0xCC, 0xFF))
add_rect(slide, 0.7, 3.3, 5, 0.05, AZUL_MEDIO)
add_textbox(slide, "Cliente: Cadena de Supermercados  |  3 Sucursales  |  3 Meses de datos",
            0.7, 3.5, 8.8, 0.45, font_size=13, color=RGBColor(0xCC, 0xDD, 0xEE))

kpi_box(slide, 0.7, 4.4, 2.5, 1.3, "1,000", "Transacciones", AZUL_MEDIO)
kpi_box(slide, 3.4, 4.4, 2.5, 1.3, "$322,967", "Ventas Totales", RGBColor(0x1A, 0x8A, 0x5A))
kpi_box(slide, 6.1, 4.4, 2.5, 1.3, "$322.97", "Ticket Promedio", RGBColor(0xC0, 0x39, 0x2B))

add_textbox(slide, "Febrero 2026", 0.7, 6.9, 9, 0.4,
            font_size=10, color=RGBColor(0x88, 0x99, 0xAA))

# ============================================================
# SLIDE 2 — CONTEXTO
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, GRIS_CLARO)
slide_header(slide, "Contexto del Análisis", "¿Quién es el cliente y qué necesita?")
slide_footer(slide)

add_rect(slide, 0.3, 1.4, 4.3, 0.42, AZUL_MEDIO)
add_textbox(slide, "¿QUIÉN?", 0.3, 1.4, 4.3, 0.42,
            font_size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

bullets_quien = [
    "Cadena de supermercados con 3 sucursales",
    "Sucursal A — Yangon:       340 transacciones",
    "Sucursal B — Mandalay:   332 transacciones",
    "Sucursal C — Naypyitaw:  328 transacciones",
    "Horario: 10:00 a 21:00 hrs (12 hs diarias)",
    "Pagos: Ewallet, Cash, Credit card",
]
y = 1.9
for b in bullets_quien:
    add_textbox(slide, f"▸  {b}", 0.4, y, 4.1, 0.37, font_size=11, color=GRIS_TEXTO)
    y += 0.37

add_rect(slide, 5.1, 1.4, 4.5, 0.42, ROJO_ACENTO)
add_textbox(slide, "¿QUÉ NECESITA?", 5.1, 1.4, 4.5, 0.42,
            font_size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
add_textbox(slide,
    "Realizar 3 campañas de Marketing mensuales para maximizar los ingresos por ventas.",
    5.2, 1.9, 4.3, 0.65, font_size=11, color=GRIS_TEXTO)

add_rect(slide, 5.1, 2.65, 4.5, 0.38, AZUL_OSCURO)
add_textbox(slide, "Preguntas clave del cliente:", 5.2, 2.65, 4.3, 0.38,
            font_size=11, bold=True, color=BLANCO)

preguntas = [
    "¿Cómo compran Mujeres vs Hombres?",
    "¿Cómo se distribuyen las compras por hora?",
    "¿Cuáles son los ingresos por línea de producto?",
    "¿Ingresos por producto y género?",
    "¿Cómo son las ventas por sucursal?",
    "¿Compras por línea de productos?",
]
y = 3.1
for p in preguntas:
    add_textbox(slide, f"▸  {p}", 5.2, y, 4.3, 0.34, font_size=11, color=GRIS_TEXTO)
    y += 0.34

# ============================================================
# SLIDE 3 — COMPRAS POR GÉNERO
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, GRIS_CLARO)
slide_header(slide, "Compras por Género", "¿Cómo se comparan las compras de Mujeres vs Hombres?")
slide_footer(slide)

female_sales = df[df['Gender'] == 'Female']['Total'].sum()
male_sales   = df[df['Gender'] == 'Male']['Total'].sum()
total_sales  = female_sales + male_sales
female_avg   = df[df['Gender'] == 'Female']['Total'].mean()
male_avg     = df[df['Gender'] == 'Male']['Total'].mean()
female_count = len(df[df['Gender'] == 'Female'])
male_count   = len(df[df['Gender'] == 'Male'])

kpi_box(slide, 0.3, 1.4, 2.9, 1.1,
        f"${female_sales:,.0f}",
        f"Ventas Mujeres ({female_sales/total_sales*100:.1f}%)",
        RGBColor(0xE9, 0x1E, 0x8C))
kpi_box(slide, 3.4, 1.4, 2.9, 1.1,
        f"${male_sales:,.0f}",
        f"Ventas Hombres ({male_sales/total_sales*100:.1f}%)",
        AZUL_MEDIO)
kpi_box(slide, 6.5, 1.4, 3.1, 1.1,
        f"{female_count} / {male_count}",
        "Transacciones Mujeres / Hombres",
        AZUL_OSCURO)

add_rect(slide, 0.3, 2.65, 9.2, 0.62, RGBColor(0xFF, 0xF3, 0xCD))
add_textbox(slide,
    f"💡  Las mujeres generan el {female_sales/total_sales*100:.1f}% de las ventas con ticket promedio "
    f"${female_avg:.2f} vs ${male_avg:.2f} de los hombres.",
    0.45, 2.67, 9.0, 0.58, font_size=13, bold=True, color=RGBColor(0x85, 0x65, 0x04))

add_placeholder_box(slide, 0.3, 3.4, 9.2, 3.55,
    "[ Insertar gráfico de barras comparativo: Ventas totales y ticket promedio por Género ]")

# ============================================================
# SLIDE 4 — DISTRIBUCIÓN POR HORAS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, GRIS_CLARO)
slide_header(slide, "Distribución de Compras por Hora", "¿En qué momentos del día se concentran las ventas?")
slide_footer(slide)

morning   = df[df['Hour'].between(10, 13)]['Total'].sum()
afternoon = df[df['Hour'].between(14, 17)]['Total'].sum()
evening   = df[df['Hour'].between(18, 21)]['Total'].sum()
hourly_sales = df.groupby('Hour')['Total'].sum()
peak_hour    = hourly_sales.idxmax()
peak_sales   = hourly_sales.max()

kpi_box(slide, 0.3, 1.4, 2.9, 1.0, f"${morning:,.0f}",
        f"Mañana 10-13h  ({morning/total_sales*100:.1f}%)", AZUL_MEDIO)
kpi_box(slide, 3.4, 1.4, 2.9, 1.0, f"${afternoon:,.0f}",
        f"Tarde 14-17h  ({afternoon/total_sales*100:.1f}%)", AZUL_OSCURO)
kpi_box(slide, 6.5, 1.4, 3.1, 1.0, f"${evening:,.0f}",
        f"Noche 18-21h  ({evening/total_sales*100:.1f}%)", RGBColor(0x6C, 0x3A, 0x83))

add_rect(slide, 0.3, 2.55, 9.2, 0.62, RGBColor(0xFF, 0xF3, 0xCD))
add_textbox(slide,
    f"💡  Hora pico: {peak_hour}:00 hs con ${peak_sales:,.0f} en ventas  |  "
    "Distribución equilibrada entre los 3 períodos del día.",
    0.45, 2.57, 9.0, 0.58, font_size=13, bold=True, color=RGBColor(0x85, 0x65, 0x04))

add_placeholder_box(slide, 0.3, 3.3, 9.2, 3.65,
    "[ Insertar gráfico de líneas: Ventas por Hora del Día (10h a 21h) ]")

# ============================================================
# SLIDE 5 — INGRESOS POR LÍNEA DE PRODUCTOS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, GRIS_CLARO)
slide_header(slide, "Ingresos por Línea de Productos", "¿Qué categorías generan más ingresos?")
slide_footer(slide)

product_sales = df.groupby('Product line')['Total'].sum().sort_values(ascending=False)

kpi_box(slide, 0.3, 1.4, 4.5, 1.0,
        product_sales.index[0], f"Línea líder  ${product_sales.iloc[0]:,.0f}", VERDE_ACENTO)
kpi_box(slide, 5.1, 1.4, 4.5, 1.0,
        product_sales.index[-1], f"Menor  ${product_sales.iloc[-1]:,.0f}", ROJO_ACENTO)

add_rect(slide, 0.3, 2.55, 9.2, 0.4, AZUL_OSCURO)
add_textbox(slide, "Línea de Producto",  0.4, 2.57, 4.5, 0.35, font_size=11, bold=True, color=BLANCO)
add_textbox(slide, "Ventas Totales",     5.2, 2.57, 2.2, 0.35, font_size=11, bold=True, color=BLANCO)
add_textbox(slide, "% del Total",        7.5, 2.57, 1.8, 0.35, font_size=11, bold=True, color=BLANCO)

y = 3.0
for i, (product, sales) in enumerate(product_sales.items()):
    pct = sales / df['Total'].sum() * 100
    bg = RGBColor(0xE8, 0xF4, 0xFD) if i % 2 == 0 else BLANCO
    add_rect(slide, 0.3, y, 9.2, 0.38, bg)
    add_textbox(slide, product,          0.4, y+0.02, 4.5, 0.35, font_size=11, color=GRIS_TEXTO)
    add_textbox(slide, f"${sales:,.2f}", 5.2, y+0.02, 2.2, 0.35, font_size=11, color=GRIS_TEXTO)
    add_textbox(slide, f"{pct:.1f}%",    7.5, y+0.02, 1.8, 0.35, font_size=11, color=GRIS_TEXTO)
    y += 0.38

add_placeholder_box(slide, 0.3, 5.35, 9.2, 1.6,
    "[ Insertar gráfico de barras horizontales: Ingresos por Línea de Producto ]")

# ============================================================
# SLIDE 6 — INGRESOS POR PRODUCTO Y GÉNERO
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, GRIS_CLARO)
slide_header(slide, "Ingresos por Producto y Género", "¿Qué compra cada género?")
slide_footer(slide)

gender_product = []
for product in df['Product line'].unique():
    f_s = df[(df['Gender'] == 'Female') & (df['Product line'] == product)]['Total'].sum()
    m_s = df[(df['Gender'] == 'Male')   & (df['Product line'] == product)]['Total'].sum()
    tot = f_s + m_s
    gap = abs(f_s/tot*100 - m_s/tot*100)
    gender_product.append((product, f_s/tot*100, m_s/tot*100, gap))

gender_product.sort(key=lambda x: x[3], reverse=True)

add_rect(slide, 0.3, 1.35, 9.2, 0.62, RGBColor(0xFF, 0xEB, 0xEE))
add_textbox(slide,
    "🎯  Mayor brecha: Health & Beauty — Hombres 62.3% vs Mujeres 37.7%  (brecha 24.5%)",
    0.45, 1.37, 9.0, 0.58, font_size=14, bold=True, color=ROJO_ACENTO)

add_rect(slide, 0.3, 2.1, 9.2, 0.4, AZUL_OSCURO)
add_textbox(slide, "Línea de Producto",  0.4, 2.12, 3.5, 0.35, font_size=11, bold=True, color=BLANCO)
add_textbox(slide, "% Mujeres",          4.1, 2.12, 1.8, 0.35, font_size=11, bold=True, color=BLANCO)
add_textbox(slide, "% Hombres",          6.0, 2.12, 1.8, 0.35, font_size=11, bold=True, color=BLANCO)
add_textbox(slide, "Brecha",             7.9, 2.12, 1.4, 0.35, font_size=11, bold=True, color=BLANCO)

y = 2.55
for i, (product, f_pct, m_pct, gap) in enumerate(gender_product):
    bg = RGBColor(0xFF, 0xEB, 0xEE) if i == 0 else (RGBColor(0xE8, 0xF4, 0xFD) if i % 2 == 0 else BLANCO)
    add_rect(slide, 0.3, y, 9.2, 0.38, bg)
    add_textbox(slide, product,         0.4, y+0.02, 3.5, 0.35, font_size=11, color=GRIS_TEXTO, bold=(i==0))
    add_textbox(slide, f"{f_pct:.1f}%", 4.1, y+0.02, 1.8, 0.35, font_size=11, color=GRIS_TEXTO)
    add_textbox(slide, f"{m_pct:.1f}%", 6.0, y+0.02, 1.8, 0.35, font_size=11, color=GRIS_TEXTO)
    add_textbox(slide, f"{gap:.1f}%",   7.9, y+0.02, 1.4, 0.35, font_size=11,
                color=ROJO_ACENTO if i == 0 else GRIS_TEXTO, bold=(i==0))
    y += 0.38

add_placeholder_box(slide, 0.3, 5.3, 9.2, 1.65,
    "[ Insertar gráfico de barras apiladas: % por Género en cada Línea de Producto ]")

# ============================================================
# SLIDE 7 — COMPRAS POR SUCURSAL
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, GRIS_CLARO)
slide_header(slide, "Compras por Sucursal", "¿Cómo se desempeña cada sucursal?")
slide_footer(slide)

branch_info   = {'A': 'Yangon', 'B': 'Mandalay', 'C': 'Naypyitaw'}
branch_colors = [AZUL_MEDIO, AZUL_OSCURO, VERDE_ACENTO]
branch_sales  = df.groupby('Branch')['Total'].sum().sort_values(ascending=False)

x = 0.3
for i, (branch, sales) in enumerate(branch_sales.items()):
    city  = branch_info[branch]
    count = len(df[df['Branch'] == branch])
    pct   = sales / df['Total'].sum() * 100
    kpi_box(slide, x, 1.4, 3.0, 1.4,
            f"${sales:,.0f}",
            f"Suc. {branch} · {city}\n{count} trans. · {pct:.1f}%",
            branch_colors[i])
    x += 3.2

add_rect(slide, 0.3, 3.0, 9.2, 0.4, AZUL_OSCURO)
add_textbox(slide, "Ticket Promedio por Sucursal:", 0.4, 3.02, 9.0, 0.35,
            font_size=12, bold=True, color=BLANCO)

x = 0.3
for i, branch in enumerate(['C', 'A', 'B']):
    avg = df[df['Branch'] == branch]['Total'].mean()
    add_rect(slide, x, 3.45, 3.0, 0.7, branch_colors[i])
    add_textbox(slide, f"Suc. {branch}: ${avg:.2f}", x, 3.47, 3.0, 0.65,
                font_size=14, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    x += 3.2

add_placeholder_box(slide, 0.3, 4.3, 9.2, 2.65,
    "[ Insertar gráfico de barras agrupadas: Ventas y Transacciones por Sucursal ]")

# Guardar parte 1
prs.save('/Users/javierfernandez/Desktop/Taller/Desafio3_Presentacion.pptx')
print(f"Slides 1-7 creadas. Presentación guardada ({len(prs.slides)} slides).")
