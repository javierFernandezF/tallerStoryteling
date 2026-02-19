from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd

# Cargar datos
def load_data():
    data_rows = []
    with open('/Users/javierfernandez/Desktop/Taller/Libro1.csv', 'r', encoding='utf-8-sig') as f:
        lines = f.readlines()
    header_line = lines[0].strip()
    if header_line.startswith('"') and header_line.endswith('"'):
        header = header_line[1:-1].split(';')
    for line in lines[1:]:
        line = line.strip()
        if line.startswith('"') and line.endswith('"'):
            row_data = line[1:-1].split(';')
            data_rows.append(row_data)
    df = pd.DataFrame(data_rows, columns=header)
    numeric_cols = ['Unit price', 'Quantity', 'Tax 5%', 'Total', 'Time', 'cogs',
                    'gross margin percentage', 'gross income', 'Rating']
    for col in numeric_cols:
        if col in df.columns:
            df[col] = df[col].str.replace(',', '.').astype(float)
    df['Hour'] = (df['Time'] * 24).round().astype(int).clip(0, 23)
    return df

df = load_data()

AZUL_OSCURO  = RGBColor(0x1A, 0x37, 0x5E)
AZUL_MEDIO   = RGBColor(0x27, 0x6F, 0xBF)
VERDE_ACENTO = RGBColor(0x2E, 0xCC, 0x71)
ROJO_ACENTO  = RGBColor(0xE7, 0x4C, 0x3C)
GRIS_CLARO   = RGBColor(0xF4, 0xF6, 0xF9)
GRIS_TEXTO   = RGBColor(0x55, 0x55, 0x55)
BLANCO       = RGBColor(0xFF, 0xFF, 0xFF)
AMARILLO     = RGBColor(0xFF, 0xF3, 0xCD)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h, font_size=12, bold=False,
                color=None, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox

def add_placeholder_box(slide, l, t, w, h, label="[ Insertar gráfico de Excel aquí ]"):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEC, 0xF0, 0xF1)
    shape.line.color.rgb = AZUL_MEDIO
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(13)
    run.font.color.rgb = AZUL_MEDIO
    run.font.italic = True
    return shape

def kpi_box(slide, l, t, w, h, value, label, color):
    add_rect(slide, l, t, w, h, color)
    add_textbox(slide, value, l, t + 0.08, w, h * 0.55,
                font_size=24, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, l, t + h * 0.55, w, h * 0.42,
                font_size=10, color=BLANCO, align=PP_ALIGN.CENTER)

def slide_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 10, 1.2, AZUL_OSCURO)
    add_textbox(slide, title, 0.3, 0.15, 9.4, 0.7,
                font_size=24, bold=True, color=BLANCO)
    if subtitle:
        add_textbox(slide, subtitle, 0.3, 0.78, 9.4, 0.35,
                    font_size=12, color=RGBColor(0xAA, 0xCC, 0xFF))

def slide_footer(slide):
    add_rect(slide, 0, 7.28, 10, 0.22, AZUL_OSCURO)
    add_textbox(slide, "Desafío 3  ·  Supermarket Sales Analysis  ·  2026",
                0.2, 7.28, 9.6, 0.22, font_size=9,
                color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.RIGHT)

# Abrir presentación existente
prs = Presentation('/Users/javierfernandez/Desktop/Taller/Desafio3_Presentacion.pptx')
blank_layout = prs.slide_layouts[6]

# ============================================================
# SLIDE 8 — COMPRAS POR LÍNEA (TRANSACCIONES)
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, GRIS_CLARO)
slide_header(slide, "Compras por Línea de Productos", "¿Qué categorías tienen más transacciones?")
slide_footer(slide)

transaction_rank = df.groupby('Product line').size().sort_values(ascending=False)

kpi_box(slide, 0.3, 1.4, 4.5, 1.0,
        transaction_rank.index[0],
        f"Más transacciones: {transaction_rank.iloc[0]}", VERDE_ACENTO)
kpi_box(slide, 5.1, 1.4, 4.5, 1.0,
        transaction_rank.index[-1],
        f"Menos transacciones: {transaction_rank.iloc[-1]}", ROJO_ACENTO)

add_rect(slide, 0.3, 2.55, 9.2, 0.4, AZUL_OSCURO)
add_textbox(slide, "#",                 0.4,  2.57, 0.5, 0.35, font_size=11, bold=True, color=BLANCO)
add_textbox(slide, "Línea de Producto", 0.9,  2.57, 4.0, 0.35, font_size=11, bold=True, color=BLANCO)
add_textbox(slide, "Transacciones",     5.1,  2.57, 2.0, 0.35, font_size=11, bold=True, color=BLANCO)
add_textbox(slide, "Ticket Promedio",   7.2,  2.57, 2.1, 0.35, font_size=11, bold=True, color=BLANCO)

y = 3.0
for i, (product, count) in enumerate(transaction_rank.items()):
    avg_ticket = df[df['Product line'] == product]['Total'].mean()
    bg = RGBColor(0xE8, 0xF8, 0xF1) if i == 0 else (RGBColor(0xE8, 0xF4, 0xFD) if i % 2 == 0 else BLANCO)
    add_rect(slide, 0.3, y, 9.2, 0.38, bg)
    add_textbox(slide, str(i+1),           0.4, y+0.02, 0.5, 0.35, font_size=11, color=GRIS_TEXTO, bold=(i==0))
    add_textbox(slide, product,            0.9, y+0.02, 4.0, 0.35, font_size=11, color=GRIS_TEXTO, bold=(i==0))
    add_textbox(slide, str(count),         5.1, y+0.02, 2.0, 0.35, font_size=11, color=GRIS_TEXTO)
    add_textbox(slide, f"${avg_ticket:.2f}", 7.2, y+0.02, 2.1, 0.35, font_size=11, color=GRIS_TEXTO)
    y += 0.38

add_placeholder_box(slide, 0.3, 5.35, 9.2, 1.6,
    "[ Insertar gráfico de barras horizontales: Transacciones por Línea de Producto ]")

# ============================================================
# SLIDE 9 — 3 OPORTUNIDADES (FOCO)
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, GRIS_CLARO)
slide_header(slide, "3 Oportunidades para Aumentar las Ventas",
             "Insights clave para las 3 campañas de marketing")
slide_footer(slide)

# Oportunidad 1
add_rect(slide, 0.2, 1.35, 9.3, 0.42, ROJO_ACENTO)
add_textbox(slide, "OPORTUNIDAD 1 — Desequilibrio de Género en Productos",
            0.3, 1.37, 9.1, 0.38, font_size=13, bold=True, color=BLANCO)

add_rect(slide, 0.2, 1.77, 9.3, 0.9, RGBColor(0xFF, 0xEB, 0xEE))
add_textbox(slide,
    "Health & Beauty: Hombres 62.3% vs Mujeres 37.7% (brecha 24.5%)  |  "
    "Food & Beverages: Mujeres 59.1% vs Hombres 40.9% (brecha 18.2%)\n"
    "→ El género minoritario en cada producto está subrepresentado. "
    "Campaña: atraer al género minoritario con comunicación y promociones dirigidas.",
    0.35, 1.79, 9.0, 0.86, font_size=11, color=GRIS_TEXTO)

# Oportunidad 2
add_rect(slide, 0.2, 2.8, 9.3, 0.42, AZUL_MEDIO)
add_textbox(slide, "OPORTUNIDAD 2 — Diferencia Member vs Clientes Normal",
            0.3, 2.82, 9.1, 0.38, font_size=13, bold=True, color=BLANCO)

member_avg = df[df['Customer type'] == 'Member']['Total'].mean()
normal_avg = df[df['Customer type'] == 'Normal']['Total'].mean()
diff = member_avg - normal_avg
normal_count = len(df[df['Customer type'] == 'Normal'])

add_rect(slide, 0.2, 3.22, 9.3, 0.9, RGBColor(0xEB, 0xF5, 0xFB))
add_textbox(slide,
    f"Members: ticket promedio ${member_avg:.2f}  |  Normal: ${normal_avg:.2f}  |  "
    f"Diferencia: ${diff:.2f} más por transacción  |  {normal_count} clientes Normal activos\n"
    "→ En Health & Beauty los Members gastan 19.7% más que Normal. "
    "Campaña: programa de fidelización para convertir clientes Normal a Member.",
    0.35, 3.24, 9.0, 0.86, font_size=11, color=GRIS_TEXTO)

# Oportunidad 3
add_rect(slide, 0.2, 4.25, 9.3, 0.42, RGBColor(0x6C, 0x3A, 0x83))
add_textbox(slide, "OPORTUNIDAD 3 — Horarios con Baja Densidad de Ventas",
            0.3, 4.27, 9.1, 0.38, font_size=13, bold=True, color=BLANCO)

avg_hourly = df.groupby('Hour')['Total'].sum().mean()
h10 = df.groupby('Hour')['Total'].sum()[10]
h21 = df.groupby('Hour')['Total'].sum()[21]

add_rect(slide, 0.2, 4.67, 9.3, 0.9, RGBColor(0xF5, 0xEE, 0xF8))
add_textbox(slide,
    f"Promedio por hora: ${avg_hourly:,.0f}  |  10:00h: ${h10:,.0f} (51 trans.)  |  "
    f"21:00h: ${h21:,.0f} (38 trans.)\n"
    "→ Las horas de apertura y cierre están por debajo del promedio. "
    "Campaña: promociones especiales en horarios de apertura y cierre.",
    0.35, 4.69, 9.0, 0.86, font_size=11, color=GRIS_TEXTO)

# Llamado a la acción
add_rect(slide, 0.2, 5.7, 9.3, 0.55, AZUL_OSCURO)
add_textbox(slide,
    "💡  Estas 3 oportunidades están basadas en patrones reales del comportamiento actual de los clientes.",
    0.35, 5.72, 9.0, 0.5, font_size=12, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 10 — CIERRE
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, AZUL_OSCURO)

add_rect(slide, 0, 0, 0.35, 7.5, AZUL_MEDIO)

add_textbox(slide, "Conclusiones", 0.7, 1.2, 8.8, 0.9,
            font_size=38, bold=True, color=BLANCO)
add_rect(slide, 0.7, 2.1, 5, 0.05, AZUL_MEDIO)

conclusiones = [
    "✓  La distribución de ventas es equilibrada entre géneros, sucursales y horarios.",
    "✓  Existen brechas de género claras en productos específicos (Health & Beauty, Food).",
    "✓  Los clientes Member generan mayor ticket promedio que los Normal.",
    "✓  Los horarios de apertura (10h) y cierre (21h) tienen potencial de crecimiento.",
    "✓  Food & Beverages lidera en ingresos; Fashion Accessories en transacciones.",
]
y = 2.3
for c in conclusiones:
    add_textbox(slide, c, 0.7, y, 8.8, 0.45, font_size=13, color=RGBColor(0xCC, 0xDD, 0xFF))
    y += 0.48

add_textbox(slide, "¿Preguntas?", 0.7, 5.5, 8.8, 0.7,
            font_size=28, bold=True, color=AZUL_MEDIO, align=PP_ALIGN.CENTER)
add_textbox(slide, "Desafío 3  ·  Supermarket Sales Analysis  ·  2026",
            0.7, 6.2, 8.8, 0.4, font_size=11,
            color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)

prs.save('/Users/javierfernandez/Desktop/Taller/Desafio3_Presentacion.pptx')
print("Slides 8-10 agregadas. Presentación completa guardada.")
